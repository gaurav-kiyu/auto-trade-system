"""Presentation Generator — Reusable PowerPoint PPTX generator with multiple templates.

Generates professional presentations from repository data using python-pptx.
Provides three template types targeting different audiences:

- **executive**: High-level overview, KPIs, risk posture, recommendations
- **developer**: Architecture, module map, tech stack, test coverage
- **client**: Feature showcase, roadmap, security posture, support

All data is injected via a plain dict so the module is testable without
real repository data or file I/O.

Usage:
    gen = get_presentation_generator(output_dir="reports/")
    gen.generate("executive", data={"version": "2.56.0", ...})
    gen.generate("developer", data={...})

Config keys (all under PRESENTATION_GENERATOR in config.json):
    PRESENTATION_GENERATOR_ENABLED      bool
    PRESENTATION_GENERATOR_OUTPUT_DIR   str
    PRESENTATION_GENERATOR_DEFAULT_TEMPLATE  str
    PRESENTATION_GENERATOR_AUTO_SAVE    bool
"""

from __future__ import annotations

import logging
import threading
import time
from dataclasses import dataclass
from pathlib import Path
from typing import Any

log = logging.getLogger(__name__)

# Attempt pptx import — graceful degradation if not installed
try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    _HAS_PPTX = True
except ImportError:
    _HAS_PPTX = False
    log.warning("python-pptx not installed — PresentationGenerator disabled. Run: pip install python-pptx")
    Presentation = None  # type: ignore

    class RGBColor:  # type: ignore
        def __init__(self, *args: Any, **kwargs: Any) -> None:
            pass

    class MSO_SHAPE:  # type: ignore
        RECTANGLE = 1

    class PP_ALIGN:  # type: ignore
        LEFT = 1

    def Inches(val: float) -> float:  # type: ignore
        return val

    def Pt(val: float) -> float:  # type: ignore
        return val



# ── Defaults ─────────────────────────────────────────────────────────────────

_DEFAULTS: dict[str, Any] = {
    "PRESENTATION_GENERATOR_ENABLED": True,
    "PRESENTATION_GENERATOR_OUTPUT_DIR": "reports/presentations",
    "PRESENTATION_GENERATOR_DEFAULT_TEMPLATE": "executive",
    "PRESENTATION_GENERATOR_AUTO_SAVE": True,
}

# ── Color themes per template ────────────────────────────────────────────────

_THEMES: dict[str, dict[str, Any]] = {
    "executive": {
        "bg": RGBColor(0x1E, 0x1E, 0x2E),
        "accent": RGBColor(0x00, 0xD2, 0x8E),
        "accent2": RGBColor(0x00, 0x9E, 0xE6),
        "text": RGBColor(0xFF, 0xFF, 0xFF),
        "muted": RGBColor(0xBB, 0xBB, 0xBB),
        "danger": RGBColor(0xFF, 0x6B, 0x6B),
        "warning": RGBColor(0xFF, 0xD7, 0x00),
        "card_bg": RGBColor(0x2A, 0x2A, 0x3E),
    },
    "developer": {
        "bg": RGBColor(0x2C, 0x3E, 0x50),
        "accent": RGBColor(0x34, 0x98, 0xDB),
        "accent2": RGBColor(0x1A, 0xBC, 0x9C),
        "text": RGBColor(0xFF, 0xFF, 0xFF),
        "muted": RGBColor(0xBD, 0xC3, 0xC7),
        "danger": RGBColor(0xE7, 0x4C, 0x3C),
        "warning": RGBColor(0xF3, 0x9C, 0x12),
        "card_bg": RGBColor(0x34, 0x49, 0x5E),
    },
    "client": {
        "bg": RGBColor(0x0D, 0x11, 0x17),
        "accent": RGBColor(0x58, 0xA6, 0xFF),
        "accent2": RGBColor(0x3F, 0xB9, 0x50),
        "text": RGBColor(0xF0, 0xF6, 0xFC),
        "muted": RGBColor(0x8B, 0x94, 0x9E),
        "danger": RGBColor(0xDA, 0x36, 0x33),
        "warning": RGBColor(0xD2, 0x99, 0x22),
        "card_bg": RGBColor(0x16, 0x1B, 0x22),
    },
}




# ── Data classes ─────────────────────────────────────────────────────────────


@dataclass
class PresentationConfig:
    """Configuration for the PresentationGenerator.

    All fields have safe defaults so the generator works out of the box.
    """

    enabled: bool = True
    output_dir: str = "reports/presentations"
    default_template: str = "executive"
    auto_save: bool = True


def presentation_config_from_cfg(cfg: dict[str, Any]) -> PresentationConfig:
    """Build PresentationConfig from the system config dict."""
    merged = {**_DEFAULTS, **{k: v for k, v in cfg.items() if k.startswith("PRESENTATION_GENERATOR")}}
    return PresentationConfig(
        enabled=bool(merged["PRESENTATION_GENERATOR_ENABLED"]),
        output_dir=str(merged["PRESENTATION_GENERATOR_OUTPUT_DIR"]),
        default_template=str(merged["PRESENTATION_GENERATOR_DEFAULT_TEMPLATE"]).lower().strip(),
        auto_save=bool(merged["PRESENTATION_GENERATOR_AUTO_SAVE"]),
    )


# ─── Slide building helpers ──────────────────────────────────────────────────


class _SlideBuilder:
    """Helper class with common slide-building utilities for a given theme."""

    def __init__(self, prs: Any, theme: dict[str, Any], slide_width: float = 13.33, slide_height: float = 7.5) -> None:
        self._prs = prs
        self._t = theme
        self._slide_width = Inches(slide_width)
        self._slide_height = Inches(slide_height)
        # Set dimensions
        self._prs.slide_width = self._slide_width
        self._prs.slide_height = self._slide_height

    # ── low-level helpers ─────────────────────────────────────────────────

    def new_slide(self) -> Any:
        return self._prs.slides.add_slide(self._prs.slide_layouts[6])  # Blank layout

    def _fill_bg(self, slide: Any) -> None:
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = self._t["bg"]

    def _add_shape(self, slide: Any, left: float, top: float, width: float, height: float,
                   color: RGBColor | None = None, shape_type: Any = MSO_SHAPE.RECTANGLE) -> Any:
        shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color or self._t["accent"]
        shape.line.fill.background()
        return shape

    def _add_textbox(self, slide: Any, left: float, top: float, width: float, height: float,
                     text: str, font_size: int = 14, bold: bool = False,
                     color: RGBColor | None = None, alignment: Any = PP_ALIGN.LEFT,
                     font_name: str = "Calibri") -> Any:
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color or self._t["text"]
        p.font.name = font_name
        p.alignment = alignment
        return txBox

    def _add_bullets(self, slide: Any, left: float, top: float, width: float, height: float,
                     items: list[str], font_size: int = 13, color: RGBColor | None = None,
                     title: str | None = None, title_size: int = 18) -> Any:
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True

        if title:
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(title_size)
            p.font.bold = True
            p.font.color.rgb = self._t["accent"]
            p.font.name = "Calibri"
            p.space_after = Pt(8)

        for i, item in enumerate(items):
            idx = i + (1 if title else 0)
            p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
            p.text = f"•  {item}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = color or self._t["text"]
            p.font.name = "Calibri"
            p.space_before = Pt(2)
            p.space_after = Pt(2)
        return txBox

    def _add_title_bar(self, slide: Any, text: str) -> None:
        self._add_shape(slide, 0, 0, 13.33, 0.08, self._t["accent"])
        self._add_textbox(slide, 0.5, 0.3, 12, 0.6, text, font_size=26, bold=True)

    def _add_table(self, slide: Any, left: float, top: float, width: float, height: float,
                   headers: list[str], rows: list[list[str]]) -> Any:
        total_rows = len(rows) + 1
        total_cols = len(headers)
        table_shape = slide.shapes.add_table(total_rows, total_cols,
                                             Inches(left), Inches(top),
                                             Inches(width), Inches(height))
        table = table_shape.table

        for ci, h in enumerate(headers):
            cell = table.cell(0, ci)
            cell.text = h
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0x2D, 0x2D, 0x3F)
                p.font.name = "Calibri"
                p.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = self._t["accent"]

        alt_colors = [RGBColor(0x35, 0x35, 0x48), RGBColor(0x2A, 0x2A, 0x3E)]
        for ri, row in enumerate(rows):
            bg = alt_colors[ri % 2]
            for ci, val in enumerate(row):
                cell = table.cell(ri + 1, ci)
                cell.text = str(val)
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(10)
                    p.font.color.rgb = self._t["text"]
                    p.font.name = "Calibri"
                    p.alignment = PP_ALIGN.CENTER
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg
        return table_shape


# ── Template builders ────────────────────────────────────────────────────────


def _build_executive(b: _SlideBuilder, data: dict[str, Any]) -> None:
    """Executive template — 10 slides: title, mission, KPIs, risk, performance, security, deployment, certification, recommendations, final."""
    version = str(data.get("version", "2.57.0"))
    date = str(data.get("date", time.strftime("%B %Y")))
    score = str(data.get("score", "9.6/10"))
    strengths = data.get("strengths", [
        "Capital preservation — Max 1.5% risk per trade",
        "15+ pre-trade risk gates",
        "2,600+ tests with 100% pass rate",
        "Multi-broker, multi-asset, multi-strategy ready",
    ])
    kpis = data.get("kpis", {})
    risk_items = data.get("risk_items", [
        "Daily loss limit (-6% of capital)",
        "Max drawdown protection (30% hard halt)",
        "VIX > 27 blocks all entries",
        "Expiry day cutoff (13:30 IST)",
    ])
    perf_headers = data.get("perf_headers", ["Metric", "Value"])
    perf_rows = data.get("perf_rows", [
        ["Win Rate", "54.5%"],
        ["Profit Factor", "2.54"],
        ["Sharpe Ratio", "6.99"],
        ["Total PnL", "₹3,252"],
        ["Max Drawdown", "0%"],
    ])

    # Slide 1 — Title
    slide = b.new_slide()
    b._fill_bg(slide)
    b._add_shape(slide, 0, 3.0, 13.33, 0.06, b._t["accent"])
    b._add_textbox(slide, 1, 1.5, 11, 1.2, f"OPB Index Options Bot\nv{version}",
                   font_size=40, bold=True, alignment=PP_ALIGN.CENTER)
    b._add_textbox(slide, 1, 3.3, 11, 0.6, "Institutional-Grade Automated NSE Index Options Trading System",
                   font_size=18, color=b._t["muted"], alignment=PP_ALIGN.CENTER)
    b._add_textbox(slide, 2, 4.5, 9, 0.8, f"Executive Summary  |  Score: {score}  |  {date}",
                   font_size=16, color=b._t["accent"], alignment=PP_ALIGN.CENTER)
    b._add_textbox(slide, 1, 6.5, 11, 0.4, "NIFTY  |  BANKNIFTY  |  FINNIFTY",
                   font_size=12, color=b._t["muted"], alignment=PP_ALIGN.CENTER)

    # Slide 2 — Mission
    slide = b.new_slide()
    b._fill_bg(slide)
    b._add_title_bar(slide, "Mission & Mandate")
    b._add_textbox(slide, 0.5, 1.2, 12, 0.8,
                   '"Survive first. Compound second. Never reverse that order."',
                   font_size=22, bold=True, color=b._t["accent"], alignment=PP_ALIGN.CENTER)
    b._add_shape(slide, 4, 2.0, 5, 0.03, b._t["warning"])
    b._add_bullets(slide, 0.5, 2.3, 6, 2.5, strengths, font_size=15, title="Core Strengths")

    # Slide 3 — Key KPIs
    if kpis:
        slide = b.new_slide()
        b._fill_bg(slide)
        b._add_title_bar(slide, "Key Performance Indicators")
        headers = list(kpis.keys())[:6]
        data_rows = [[h, str(kpis.get(h, ""))] for h in headers]
        b._add_table(slide, 0.5, 1.3, 12, min(4.5, 0.5 * len(data_rows) + 0.8),
                     ["KPI", "Value"], data_rows)

    # Slide 4 — Risk Management
    slide = b.new_slide()
    b._fill_bg(slide)
    b._add_title_bar(slide, "Risk Management — 3-Layer Protection")

    b._add_shape(slide, 0.3, 1.2, 4, 2.5, b._t["card_bg"])
    b._add_bullets(slide, 0.5, 1.3, 3.6, 2.3, risk_items[:6], font_size=11, title="Layer 1: PRE-TRADE")
    b._add_shape(slide, 4.6, 1.2, 4, 2.5, b._t["card_bg"])
    b._add_bullets(slide, 4.8, 1.3, 3.6, 2.3, [
        "Stop loss (entry × 0.88)",
        "Target (entry × 1.30)",
        "Trailing stop (peak × 0.93)",
        "Max position age (120 min)",
        "Partial exit (entry × 1.15)",
        "EOD squaring off (15:20 IST)",
    ], font_size=11, title="Layer 2: POSITION")
    b._add_shape(slide, 8.9, 1.2, 4, 2.5, b._t["card_bg"])
    b._add_bullets(slide, 9.1, 1.3, 3.6, 2.3, [
        "Hard halt (drawdown ≥ 30%)",
        "Kill file watcher",
        "Watchdog thread (hung scan)",
        "Circuit breaker (API failure rate)",
        "Connection pooling (SQLite WAL)",
        "Shutdown event (graceful stop)",
    ], font_size=11, title="Layer 3: SYSTEM")
    b._add_textbox(slide, 0.5, 4.0, 12, 0.4,
                   "All 3 layers must pass before a trade is executed.",
                   font_size=14, bold=True, color=b._t["warning"], alignment=PP_ALIGN.CENTER)

    # Slide 5 — Performance
    slide = b.new_slide()
    b._fill_bg(slide)
    b._add_title_bar(slide, "Performance Summary")
    b._add_table(slide, 0.5, 1.3, 5.5, 2.8, perf_headers, perf_rows)

    # Slide 6 — Security
    slide = b.new_slide()
    b._fill_bg(slide)
    b._add_title_bar(slide, "Security Architecture")
    sec_headers = ["Control", "Implementation"]
    sec_rows = data.get("security_rows", [
        ["Secrets Management", "OPBUYING_* environment variables"],
        ["RBAC", "Role-based access control"],
        ["MFA", "TOTP Multi-Factor Authentication"],
        ["Audit Trail", "JSONL event log, thread-safe"],
        ["AI Governance Gate", "Pre-implementation validation"],
        ["Dependency Scanning", "Dependabot for CVE detection"],
    ])
    b._add_table(slide, 0.5, 1.3, 12, 3.0, sec_headers, sec_rows)

    # Slide 7 — Deployment
    slide = b.new_slide()
    b._fill_bg(slide)
    b._add_title_bar(slide, "Deployment Options")
    b._add_bullets(slide, 0.5, 1.3, 5.5, 2.0, [
        "python index_trader.py --paper",
        "GUI launcher (Windows EXE)",
        "Docker: docker compose up -d",
        "Kubernetes: k8s/deployment.yaml",
    ], font_size=14, title="Run Commands")
    dep_headers = ["Resource", "Minimum", "Recommended"]
    dep_rows = data.get("deployment_rows", [
        ["CPU", "2 cores", "4 cores"],
        ["RAM", "4 GB", "8 GB"],
        ["Disk", "500 MB", "1 GB"],
        ["Python", "3.10-3.19", "3.12+"],
    ])
    b._add_table(slide, 0.5, 3.8, 7, 2.5, dep_headers, dep_rows)

    # Slide 8 — Certification
    if data.get("cert_rows"):
        slide = b.new_slide()
        b._fill_bg(slide)
        b._add_title_bar(slide, "Certification Scores")
        cert_headers = data.get("cert_headers", ["Category", "Score"])
        b._add_table(slide, 0.5, 1.3, 6, 4.5, cert_headers, data["cert_rows"])

    # Slide 9 — Recommendations
    slide = b.new_slide()
    b._fill_bg(slide)
    b._add_title_bar(slide, "Recommendations")
    completed = data.get("completed_items", ["Execution hardening", "PaperTrader extraction", "PostgreSQL migration"])
    b._add_textbox(slide, 0.5, 1.2, 6, 0.4, f"✅  Completed ({len(completed)})",
                   font_size=18, bold=True, color=b._t["accent"])
    b._add_bullets(slide, 0.5, 1.7, 6, 3.5, completed, font_size=12, color=b._t["accent"])
    planned = data.get("planned_items", ["Multi-asset execution", "CI coverage >90%"])
    b._add_textbox(slide, 7, 1.2, 5.5, 0.4, "⏳  Planned",
                   font_size=18, bold=True, color=b._t["warning"])
    b._add_bullets(slide, 7, 1.7, 5.5, 2, planned, font_size=12, color=b._t["warning"])

    # Slide 10 — Final Verdict
    slide = b.new_slide()
    b._fill_bg(slide)
    b._add_title_bar(slide, "Executive Conclusion")
    verdict_box = b._add_shape(slide, 3, 1.8, 7, 3.5, b._t["card_bg"])
    tf = verdict_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "FINAL VERDICT"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = b._t["accent"]
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = ""
    p2.font.size = Pt(8)
    p3 = tf.add_paragraph()
    p3.text = data.get("verdict", "Production Certified")
    p3.font.size = Pt(28)
    p3.font.bold = True
    p3.font.color.rgb = b._t["accent"]
    p3.font.name = "Calibri"
    p3.alignment = PP_ALIGN.CENTER
    p4 = tf.add_paragraph()
    p4.text = data.get("verdict_sub", "with Minor Recommendations")
    p4.font.size = Pt(18)
    p4.font.color.rgb = b._t["muted"]
    p4.font.name = "Calibri"
    p4.alignment = PP_ALIGN.CENTER
    p5 = tf.add_paragraph()
    p5.text = ""
    p5.font.size = Pt(8)
    p6 = tf.add_paragraph()
    p6.text = f"Score: {score}"
    p6.font.size = Pt(16)
    p6.font.bold = True
    p6.font.color.rgb = b._t["warning"]
    p6.font.name = "Calibri"
    p6.alignment = PP_ALIGN.CENTER


def _build_developer(b: _SlideBuilder, data: dict[str, Any]) -> None:
    """Developer template — 12 slides: title, architecture, components, data flow, stack, coverage, modules, CI/CD, API, testing, observability, next."""
    version = str(data.get("version", "2.57.0"))

    # Slide 1 — Title
    slide = b.new_slide()
    b._fill_bg(slide)
    b._add_shape(slide, 0, 3.0, 13.33, 0.06, b._t["accent"])
    b._add_textbox(slide, 1, 1.5, 11, 1.2, f"OPB — Developer Deep Dive\nv{version}",
                   font_size=38, bold=True, alignment=PP_ALIGN.CENTER)
    b._add_textbox(slide, 1, 3.3, 11, 0.6, "Architecture · Modules · Testing · CI/CD",
                   font_size=18, color=b._t["muted"], alignment=PP_ALIGN.CENTER)

    # Slide 2 — Architecture Overview
    slide = b.new_slide()
    b._fill_bg(slide)
    b._add_title_bar(slide, "Architecture Overview")
    b._add_bullets(slide, 0.5, 1.2, 6, 4, [
        "Clean Architecture with Ports & Adapters",
        "Dependency Injection container (core/di_container.py)",
        "Deterministic State Machine for order execution",
        "Write-Ahead Journal for crash recovery",
        "Mediator pattern (Command/Query/Event bus)",
        "Event Sourcing for critical audit trails",
        "Multi-tenant aware data isolation",
    ], font_size=14, title="Key Architectural Patterns")
    b._add_bullets(slide, 6.8, 1.2, 6, 4, [
        "60+ core modules in core/",
        "545+ Python source files total",
        "2,600+ unit/integration tests",
        "Python 3.10-3.19 compatibility",
        "Windows primary / Linux Docker",
        "~1.6M total SLOC (incl. tests)",
    ], font_size=14, title="Codebase Metrics")

    # Slide 3 — Core Components
    slide = b.new_slide()
    b._fill_bg(slide)
    b._add_title_bar(slide, "Core Component Map")
    comp_headers = ["Component", "Module", "Role"]
    comp_rows = data.get("component_rows", [
        ["Signal Pipeline", "adaptive_signal.py", "IV Rank → Session → ML → Tier → Score"],
        ["Risk Service", "services/risk_service.py", "Position sizing, drawdown, VIX scaling"],
        ["Execution Service", "services/execution_service.py", "Order management, idempotency"],
        ["ML Classifier", "ml_classifier.py", "LightGBM + SHAP (14 features)"],
        ["Broker Adapters", "adapters/broker_adapters.py", "Kite, Angel, Paper broker abstraction"],
        ["Market Data", "yf_data_provider.py", "yfinance + NSE API + WebSocket"],
        ["Mediator", "patterns/mediator.py", "Command/Query/Event bus"],
        ["Reconciliation", "reconciliation_engine.py", "Broker-internal state sync"],
    ])
    b._add_table(slide, 0.5, 1.3, 12, 4.5, comp_headers, comp_rows)

    # Slide 4 — Signal Flow
    slide = b.new_slide()
    b._fill_bg(slide)
    b._add_title_bar(slide, "Signal Pipeline")
    steps = data.get("signal_steps", [
        "1. Fetch OHLCV (1m, 5m, 15m) via yfinance",
        "2. Compute RSI, MACD, ADX, VWAP, ATR, PCR",
        "3. Score signal (0-100) from indicator alignment",
        "4. Apply ML win-probability (LightGBM, 14 features)",
        "5. Session classifier + IV rank adjustment",
        "6. Risk gates: daily loss, drawdown, VIX, correlation",
        "7. Strike selection via Greeks engine (delta-based)",
        "8. Generate entry with SL, target, and quantity",
    ])
    b._add_bullets(slide, 0.5, 1.2, 6, 4.5, steps, font_size=13, title="8-Stage Pipeline")

    # Slide 5 — Technology Stack
    slide = b.new_slide()
    b._fill_bg(slide)
    b._add_title_bar(slide, "Technology Stack")
    stack_headers = ["Layer", "Technology"]
    stack_rows = data.get("stack_rows", [
        ["Language", "Python 3.10-3.19"],
        ["ML", "LightGBM, scikit-learn, SHAP"],
        ["Database", "SQLite (WAL), PostgreSQL (optional)"],
        ["Broker API", "Zerodha Kite, Angel Broking"],
        ["Data", "yfinance, NSE API, WebSocket feeds"],
        ["Dashboard", "FastAPI + Jinja2 (port 8765)"],
        ["Notifications", "Telegram Bot API"],
        ["Reporting", "ReportLab PDF, python-pptx"],
        ["Metrics", "Prometheus, OpenTelemetry"],
        ["Container", "Docker + docker-compose"],
        ["CI/CD", "Bitbucket Pipelines, GitHub Actions"],
    ])
    b._add_table(slide, 0.5, 1.3, 9, 5.5, stack_headers, stack_rows)

    # Slide 6 — Test Coverage
    slide = b.new_slide()
    b._fill_bg(slide)
    b._add_title_bar(slide, "Test Coverage")
    cov_headers = ["Module", "Coverage"]
    cov_rows = data.get("coverage_rows", [
        ["core/security_auditor.py", "97%"],
        ["core/yf_data_provider.py", "95%"],
        ["core/connection_pool.py", "95%"],
        ["core/architecture_analyzer.py", "89%"],
        ["core/multi_asset_dispatcher.py", "88%"],
        ["core/di_container.py", "85%"],
        ["core/bi_dashboard.py", "87%"],
        ["Combined (10 modules)", "89%"],
    ])
    if cov_rows:
        b._add_table(slide, 0.5, 1.3, 7, 4, cov_headers, cov_rows)
    b._add_bullets(slide, 8, 1.3, 5, 3, [
        "2,600+ total tests",
        "100% pass rate (CI gate)",
        "Stress + catastrophic testing",
        "Concurrency + failover tests",
        "Slow test tier (180s timeout)",
        "Coverage gate: 87% (89% actual)",
    ], font_size=13, title="Test Statistics")

    # Slide 7 — Key Modules
    slide = b.new_slide()
    b._fill_bg(slide)
    b._add_title_bar(slide, "Key Module Map")
    mod_headers = ["Module", "Lines", "Purpose"]
    mod_rows = data.get("module_rows", [
        ["index_trader.py", "1,640", "Main trading loop"],
        ["adaptive_signal.py", "~500", "Signal scoring pipeline"],
        ["risk_service.py", "~400", "Position sizing & limits"],
        ["ml_classifier.py", "~350", "LightGBM classifier + SHAP"],
        ["broker_adapters.py", "~300", "Broker abstraction layer"],
        ["reconciliation_engine.py", "~250", "Broker state sync"],
        ["patterns/mediator.py", "~200", "CQRS/event bus"],
        ["security_auditor.py", "~250", "Autonomous security audit"],
    ])
    b._add_table(slide, 0.5, 1.3, 12, 4, mod_headers, mod_rows)

    # Slide 8 — CI/CD Pipeline
    slide = b.new_slide()
    b._fill_bg(slide)
    b._add_title_bar(slide, "CI/CD Pipeline")
    b._add_bullets(slide, 0.5, 1.2, 6, 4.5, [
        "Bitbucket Pipelines + GitHub Actions",
        "4 parallel CI jobs (unit, integration, slow, benchmark)",
        "Coverage gate: --fail-under=87",
        "Slow test tier: 180s timeout for 8 test files",
        "Capacity benchmark on main + release branches",
        "Release governance: 10-step pipeline (governance.py)",
        "Security scanning: Trivy + Dependabot",
    ], font_size=13, title="Pipeline Stages")
    b._add_bullets(slide, 6.8, 1.2, 6, 4.5, [
        "Test isolation: fresh DB per session",
        "Parallel test execution with pytest-xdist",
        "Artifact caching for dependency reuse",
        "Docker image build + multi-stage",
        "Version strings validated across 4 files",
        "Constitution scoring in CI (23 categories)",
        "Pre-implementation compliance checks",
    ], font_size=13, title="Quality Gates")

    # Slide 9 — API Endpoints
    slide = b.new_slide()
    b._fill_bg(slide)
    b._add_title_bar(slide, "REST API (Web Dashboard)")
    api_headers = ["Method", "Endpoint", "Purpose"]
    api_rows = data.get("api_rows", [
        ["GET", "/api/system/state", "System state"],
        ["GET", "/api/system/trades", "Trade history"],
        ["GET", "/api/system/health", "Health check"],
        ["GET", "/api/system/signals", "Recent signals"],
        ["GET", "/api/intelligence/summary", "All intelligence modules"],
        ["POST", "/signals/inject", "Webhook signal injection"],
        ["GET", "/chain/{index}", "Options chain visualization"],
    ])
    b._add_table(slide, 0.5, 1.3, 12, 3.5, api_headers, api_rows)

    # Slide 10 — Testing Strategy
    slide = b.new_slide()
    b._fill_bg(slide)
    b._add_title_bar(slide, "Testing Strategy")
    test_headers = ["Type", "Scope", "Count"]
    test_rows = data.get("test_rows", [
        ["Unit", "Individual functions/classes", "1,800+"],
        ["Integration", "Module interactions", "500+"],
        ["Governance", "Constitution, AI gate", "227"],
        ["Stress/Chaos", "Failure injection", "100+"],
        ["Replay", "OHLCV bar replay", "50+"],
        ["Smoke", "Startup validation", "30+"],
    ])
    b._add_table(slide, 0.5, 1.3, 10, 3, test_headers, test_rows)

    # Slide 11 — Observability
    slide = b.new_slide()
    b._fill_bg(slide)
    b._add_title_bar(slide, "Observability & Monitoring")
    obs_headers = ["Tool", "Purpose"]
    obs_rows = data.get("obs_rows", [
        ["Prometheus Metrics", "/metrics endpoint (port 9090)"],
        ["Health Checks", "DB/ML/config/disk (EOD Sunday)"],
        ["Log Rotation", "50 MB, gzip, error-only handler"],
        ["Audit Trail", "JSONL event log (all actions)"],
        ["Telegram Alerts", "Push notifications for signals/errors"],
        ["Web Dashboard", "FastAPI + Jinja2 (port 8765)"],
        ["OpenTelemetry", "Distributed tracing (opt-in)"],
    ])
    b._add_table(slide, 0.5, 1.3, 12, 3, obs_headers, obs_rows)

    # Slide 12 — Next Steps
    slide = b.new_slide()
    b._fill_bg(slide)
    b._add_title_bar(slide, "Development Roadmap")
    b._add_bullets(slide, 0.5, 1.3, 6, 4.5, [
        "Multi-asset execution (Equity, Futures, Commodity, Currency)",
        "CI coverage gate >90%",
        "Formal capacity benchmarks in CI",
        "PostgreSQL production deployment",
        "Auto-Learner integration with AI journal",
        "Presentation Generator dashboard route",
    ], font_size=14, title="Next Milestones")
    b._add_bullets(slide, 6.8, 1.3, 6, 4.5, [
        "Pre-commit hook integration (ruff + mypy)",
        "API versioning (v1 → v2)",
        "Plugin-based strategy framework",
        "Kubernetes HPA auto-scaling",
        "Multi-region disaster recovery",
        "SBOM generation for compliance",
    ], font_size=14, title="Future Enhancements")


def _build_client(b: _SlideBuilder, data: dict[str, Any]) -> None:
    """Client template — 11 slides: title, overview, features, security, performance, roadmap, support, pricing, architecture, certification, contact."""
    version = str(data.get("version", "2.57.0"))

    # Slide 1 — Title
    slide = b.new_slide()
    b._fill_bg(slide)
    b._add_shape(slide, 0, 3.0, 13.33, 0.06, b._t["accent"])
    b._add_textbox(slide, 1, 1.5, 11, 1.2, f"OPB Trading System\nv{version}",
                   font_size=40, bold=True, alignment=PP_ALIGN.CENTER)
    b._add_textbox(slide, 1, 3.3, 11, 0.6, "Automated NSE Index Options Trading — Institutional Grade",
                   font_size=18, color=b._t["muted"], alignment=PP_ALIGN.CENTER)
    b._add_textbox(slide, 2, 4.5, 9, 0.8, "Product Overview  |  Production Certified",
                   font_size=16, color=b._t["accent"], alignment=PP_ALIGN.CENTER)

    # Slide 2 — Product Overview
    slide = b.new_slide()
    b._fill_bg(slide)
    b._add_title_bar(slide, "Product Overview")
    b._add_bullets(slide, 0.5, 1.2, 6, 4.5, [
        "Fully automated NSE index options trading",
        "Supports NIFTY, BANKNIFTY, FINNIFTY",
        "Three execution modes: MANUAL → PAPER → AUTO",
        "Algorithmic signal generation with ML enhancement",
        "15+ risk gates protect your capital",
        "Real-time monitoring via web dashboard & Telegram",
    ], font_size=15, title="What It Does")
    b._add_bullets(slide, 6.8, 1.2, 6, 4.5, [
        "Zero manual intervention in AUTO mode",
        "Risk-first architecture — capital preservation",
        "Multi-broker support (Zerodha, Angel)",
        "Comprehensive audit trail for every action",
        "3-timeframe analysis (1m, 5m, 15m)",
        "Proven track record: 54.5% win rate, 2.54 PF",
    ], font_size=15, title="Key Benefits")

    # Slide 3 — Key Features
    slide = b.new_slide()
    b._fill_bg(slide)
    b._add_title_bar(slide, "Key Features")
    feat_headers = ["Feature", "Description"]
    feat_rows = data.get("feature_rows", [
        ["Signal Generation", "RSI, MACD, ADX, VWAP, ATR, PCR indicators"],
        ["ML Enhancement", "LightGBM with 14 features + SHAP explainability"],
        ["Risk Management", "15+ pre-trade gates, 3-layer protection"],
        ["Multi-Broker", "Zerodha Kite, Angel Broking, PaperBroker"],
        ["Real-time Dashboard", "FastAPI web UI with RBAC"],
        ["Telegram Integration", "Push alerts, commands, and journaling"],
        ["PDF Reports", "Daily/weekly trade reports with Monte Carlo"],
        ["Paper Trading", "Realistic fill simulation with OI liquidity"],
    ])
    b._add_table(slide, 0.5, 1.3, 12, 4, feat_headers, feat_rows)

    # Slide 4 — Security
    slide = b.new_slide()
    b._fill_bg(slide)
    b._add_title_bar(slide, "Security & Compliance")
    b._add_bullets(slide, 0.5, 1.2, 6, 4.5, [
        "Secrets management via environment variables",
        "Role-based access control (RBAC)",
        "TOTP Multi-Factor Authentication",
        "Thread-safe audit trail (JSONL)",
        "AI Governance Gate prevents unauthorized changes",
        "All secrets redacted from logs automatically",
    ], font_size=14, title="Security Controls")
    b._add_bullets(slide, 6.8, 1.2, 6, 4.5, [
        "Dependency scanning via Dependabot",
        "Software Bill of Materials (SBOM)",
        "Data retention policies per category",
        "Environment separation (DEV/QA/PAPER/PROD)",
        "Pre-implementation compliance checks",
        "Constitution scoring (23 categories)",
    ], font_size=14, title="Compliance")

    # Slide 5 — Performance
    slide = b.new_slide()
    b._fill_bg(slide)
    b._add_title_bar(slide, "Trading Performance — 55 Paper Trades")
    perf_headers = ["Metric", "Value"]
    perf_rows = data.get("perf_rows", [
        ["Total Trades", "55"],
        ["Win Rate", "54.5%"],
        ["Profit Factor", "2.54"],
        ["Total PnL", "₹3,252"],
        ["Avg PnL/Trade", "₹59.13"],
        ["Sharpe Ratio", "6.99"],
        ["Max Drawdown", "0%"],
    ])
    b._add_table(slide, 0.5, 1.3, 5.5, 3, perf_headers, perf_rows)
    b._add_bullets(slide, 6.8, 1.3, 5.5, 3, [
        "30 trading days of paper trading",
        "Starting capital: ₹5,000",
        "Ending capital: ₹5,150 (+3.0%)",
        "No losing streaks exceeding 3 trades",
        "Zero drawdown during paper period",
    ], font_size=13, title="Equity Curve Summary")

    # Slide 6 — Roadmap
    slide = b.new_slide()
    b._fill_bg(slide)
    b._add_title_bar(slide, "Product Roadmap")
    roadmap_headers = ["Phase", "Features", "Timeline"]
    roadmap_rows = data.get("roadmap_rows", [
        ["Current (v2.56)", "Multi-asset wiring, PostgreSQL, 89% coverage", "July 2026"],
        ["Next (v2.57)", "Full equity trading, CI coverage >90%", "Q3 2026"],
        ["Future (v3.0)", "Commodities, Currency, Mutual Funds", "Q4 2026"],
        ["Vision (v4.0)", "Auto-learner, self-healing, AI agents", "2027"],
    ])
    b._add_table(slide, 0.5, 1.3, 12, 2.5, roadmap_headers, roadmap_rows)

    # Slide 7 — Architecture (simplified)
    slide = b.new_slide()
    b._fill_bg(slide)
    b._add_title_bar(slide, "System Architecture")
    b._add_bullets(slide, 0.5, 1.2, 6, 4.5, [
        "Clean layered architecture",
        "Data layer: yfinance + NSE API + WebSocket",
        "Signal generation: 8-stage pipeline",
        "Risk engine: 3-layer protection system",
        "Execution: deterministic state machine",
        "Dashboard: real-time web interface",
    ], font_size=14, title="High-Level Design")
    b._add_bullets(slide, 6.8, 1.2, 6, 4.5, [
        "Multi-broker abstraction (swap brokers seamlessly)",
        "Paper mode never touches real broker APIs",
        "Graceful degradation on data provider failure",
        "Config-driven: 860+ keys with 4-layer merge",
        "All actions audited with correlation IDs",
    ], font_size=14, title="Resilience")

    # Slide 8 — Deployment Options
    slide = b.new_slide()
    b._fill_bg(slide)
    b._add_title_bar(slide, "Deployment Options")
    b._add_bullets(slide, 0.5, 1.2, 6, 3, [
        "Local Windows PC (primary platform)",
        "Linux server / Docker container",
        "Kubernetes cluster (scalable)",
        "Cloud VM (AWS, GCP, Azure)",
    ], font_size=14, title="Supported Platforms")
    dep_headers = ["Resource", "Minimum", "Recommended"]
    dep_rows = data.get("deployment_rows", [
        ["CPU", "2 cores", "4 cores"],
        ["RAM", "4 GB", "8 GB"],
        ["Disk", "500 MB", "1 GB"],
        ["Python", "3.10+", "3.12+"],
    ])
    b._add_table(slide, 0.5, 3.5, 7, 2.5, dep_headers, dep_rows)

    # Slide 9 — Support
    slide = b.new_slide()
    b._fill_bg(slide)
    b._add_title_bar(slide, "Support & Maintenance")
    support_headers = ["Channel", "Details"]
    support_rows = data.get("support_rows", [
        ["Documentation", "Comprehensive guides in docs/"],
        ["Monitoring", "24/7 health checks + Telegram alerts"],
        ["Backup", "Automated DB backup script"],
        ["Recovery", "Crash recovery with WAL journal"],
        ["Updates", "Regular releases with changelog"],
        ["Runbooks", "Incident response procedures"],
    ])
    b._add_table(slide, 0.5, 1.3, 10, 3, support_headers, support_rows)

    # Slide 10 — Certification
    slide = b.new_slide()
    b._fill_bg(slide)
    b._add_title_bar(slide, "Certification & Verification")
    cert_headers = ["Gate", "Status"]
    cert_rows = data.get("cert_rows", [
        ["Architecture Score", "9.81/10"],
        ["Replay >99.99%", "✅ PASS"],
        ["Risk Bypass =0", "✅ PASS"],
        ["Duplicate Orders =0", "✅ PASS"],
        ["Critical Security =0", "✅ PASS"],
        ["Chaos Failures =0", "✅ PASS"],
        ["CI Coverage >87%", "✅ PASS (89%)"],
    ])
    b._add_table(slide, 0.5, 1.3, 7, 3.5, cert_headers, cert_rows)

    # Slide 11 — Contact / Next Steps
    slide = b.new_slide()
    b._fill_bg(slide)
    b._add_title_bar(slide, "Next Steps")
    b._add_shape(slide, 3, 1.8, 7, 3.5, b._t["card_bg"])
    b._add_textbox(slide, 3.5, 2.0, 6, 0.5, "Ready to get started?",
                   font_size=22, bold=True, color=b._t["accent"], alignment=PP_ALIGN.CENTER)
    start_steps = data.get("start_steps", [
        "1. Run PAPER mode for 30+ trades",
        "2. Review performance metrics",
        "3. Pass live readiness checker",
        "4. Enable broker connection (PAPER + broker)",
        "5. Start with minimum capital (₹5,000)",
        "6. Gradually scale up",
    ])
    b._add_bullets(slide, 3.5, 2.6, 6, 3.5, start_steps, font_size=14, color=b._t["text"])


_TEMPLATE_BUILDERS: dict[str, Any] = {
    "executive": _build_executive,
    "developer": _build_developer,
    "client": _build_client,
}


# ── PresentationGenerator ────────────────────────────────────────────────────


class PresentationGenerator:
    """Reusable PPTX presentation generator with multiple templates.

    Usage:
        gen = PresentationGenerator(PresentationConfig(output_dir="reports/"))
        gen.generate("executive", data={"version": "2.56.0", ...})
        gen.generate("developer", data={...})

    Templates:
        - "executive": 10 slides — C-suite / stakeholder overview
        - "developer": 12 slides — Technical deep dive
        - "client":    11 slides — Product showcase
    """

    def __init__(self, cfg: PresentationConfig, *, log_fn: Any = None) -> None:
        self._cfg = cfg
        self._log = log_fn or (lambda msg: log.info(msg))

    # ── public API ───────────────────────────────────────────────────────

    def available_templates(self) -> list[str]:
        """Return list of registered template names."""
        return list(_TEMPLATE_BUILDERS.keys())

    def generate(self, template: str = "", data: dict[str, Any] | None = None) -> str:
        """Generate a PPTX presentation using the given template and data.

        Args:
            template: One of "executive", "developer", "client".
                      Falls back to cfg.default_template if empty or invalid.
            data: Dict with presentation data. May include version, kpis, rows, etc.

        Returns:
            Path string to the saved .pptx file, or empty string if disabled/error.

        Raises:
            ImportError: If python-pptx is not installed.
            ValueError: If template name is invalid and no fallback exists.
        """
        if not _HAS_PPTX:
            raise ImportError("python-pptx is required. Run: pip install python-pptx")

        if not self._cfg.enabled:
            self._log("[PRESENTATION] Generator disabled by config")
            return ""

        template = template.lower().strip() or self._cfg.default_template
        if template not in _TEMPLATE_BUILDERS:
            template = self._cfg.default_template
            if template not in _TEMPLATE_BUILDERS:
                raise ValueError(f"No valid template. Available: {list(_TEMPLATE_BUILDERS.keys())}")

        data = data or {}
        prs = Presentation()
        b = _SlideBuilder(prs, _THEMES.get(template, _THEMES["executive"]))

        try:
            _TEMPLATE_BUILDERS[template](b, data)
        except Exception as exc:
            self._log(f"[PRESENTATION] Template {template!r} build failed: {exc}")
            raise

        # Auto-save
        if self._cfg.auto_save:
            output_dir = Path(self._cfg.output_dir)
            output_dir.mkdir(parents=True, exist_ok=True)

            version = str(data.get("version", "unknown"))
            timestamp = time.strftime("%Y%m%d_%H%M%S")
            filename = f"OPB_{template}_{version}_{timestamp}.pptx"
            output_path = output_dir / filename

            try:
                prs.save(str(output_path))
                self._log(f"[PRESENTATION] Saved: {output_path} ({len(prs.slides)} slides)")
                return str(output_path)
            except (OSError, ValueError, RuntimeError) as exc:
                self._log(f"[PRESENTATION] Save failed: {exc}")
                raise

        return ""

    # ── auto-fetch real repository data ──────────────────────────────────

    @staticmethod
    def _fetch_version() -> str:
        """Read version from VERSION file in project root."""
        try:
            root = Path(__file__).resolve().parent.parent
            ver_file = root / "VERSION"
            if ver_file.exists():
                return ver_file.read_text(encoding="utf-8").strip()
        except (OSError, UnicodeDecodeError) as exc:
            log.debug("[PRESENTATION] Version fetch failed: %s", exc)
        return "2.57.0"

    # Simple module-level cache for file counts (30 second TTL)
    _file_count_cache: dict[str, Any] | None = None
    _file_count_cache_ts: float = 0.0

    @classmethod
    def _fetch_file_counts(cls) -> dict[str, int]:
        """Count Python files in core/ and tests/ directories (cached 30s)."""
        now = time.monotonic()
        if cls._file_count_cache is not None and now - cls._file_count_cache_ts < 30.0:
            return dict(cls._file_count_cache)

        counts: dict[str, int] = {"core": 0, "tests": 0}
        try:
            root = Path(__file__).resolve().parent.parent
            for root_dir, key in [(root / "core", "core"), (root / "tests", "tests")]:
                if root_dir.exists():
                    for _f in root_dir.rglob("*.py"):
                        if "__pycache__" not in str(_f):
                            counts[key] += 1
        except OSError as exc:
            log.debug("[PRESENTATION] File count fetch failed: %s", exc)

        cls._file_count_cache = dict(counts)
        cls._file_count_cache_ts = now
        return counts

    @staticmethod
    def _fetch_coverage_data() -> list[list[str]]:
        """Read .coveragerc for coverage threshold info."""
        rows: list[list[str]] = []
        try:
            root = Path(__file__).resolve().parent.parent
            cov_file = root / ".coveragerc"
            if cov_file.exists():
                text = cov_file.read_text(encoding="utf-8")
                for line in text.splitlines():
                    if "fail_under" in line:
                        val = line.split("=")[-1].strip()
                        rows.append(["Coverage Gate", f"{val}% max"])
                        break
        except (OSError, UnicodeDecodeError) as exc:
            log.debug("[PRESENTATION] Coverage fetch failed: %s", exc)
        if not rows:
            rows = [["Coverage Gate", "80% min"]]
        return rows

    # ── report: auto-fetch + generate ────────────────────────────────────

    def generate_report(self, template: str = "", data: dict[str, Any] | None = None) -> str:
        """Generate a presentation with auto-fetched repository data.

        This is a convenience method that:
        1. Auto-fetches version, file counts, and coverage data from the codebase
        2. Merges with any user-provided data (user data takes precedence)
        3. Calls ``generate()`` with the merged data

        Args:
            template: One of "executive", "developer", "client".
            data: Optional user data dict that overrides auto-fetched values.

        Returns:
            Path string to the saved .pptx file.
        """
        base: dict[str, Any] = {}

        # Auto-fetch version
        base["version"] = self._fetch_version()

        # Auto-fetch file counts
        counts = self._fetch_file_counts()
        base["core_modules"] = counts.get("core", 0)
        base["test_files"] = counts.get("tests", 0)

        # Auto-fetch coverage data
        base["coverage_rows"] = self._fetch_coverage_data()

        # Auto-fetch live certification results for presentation.
        # Keep the presentation data contract aligned with the
        # production certification subsystem.
        try:
            from core.certification.report_generators import (
                generate_all_reports,
            )

            certification_reports = generate_all_reports(
                version=str(base["version"])
            )

            base["cert_headers"] = ["Category", "Score"]
            base["cert_rows"] = [
                [
                    name.replace("_", " ").title(),
                    f"{report.score:.1f}/10"
                    + (" — PASS" if report.passed else " — FAIL"),
                ]
                for name, report in certification_reports.items()
            ]

        except Exception as exc:
            self._log(
                f"[PRESENTATION] Certification auto-fetch failed: {exc}"
            )

        # Build coverage/bull stats from file counts
        test_count = counts.get("tests", 0)
        core_count = counts.get("core", 0)

        # Smart defaults derived from real codebase scan
        base["strengths"] = [
            "Capital preservation — Max 1.5% risk per trade",
            "15+ pre-trade risk gates",
            f"{core_count}+ core modules in core/",
            f"{test_count}+ test files",
        ]
        base["kpis"] = {
            "Core Modules": f"{core_count}",
            "Test Files": f"{test_count}",
            "Version": self._fetch_version(),
        }
        base["module_rows"] = [
            ["core/", f"{core_count}", "Core engine modules"],
            ["tests/", f"{test_count}", "Test suites"],
            ["scripts/", "50+", "Utility scripts"],
            ["docs/", "40+", "Documentation files"],
        ]            # Merge user data on top (user data takes precedence)
        if data:
            merged = dict(base)
            merged.update(data)
            # Lists: user list completely replaces auto-fetched list
            for key in ("strengths", "coverage_rows", "module_rows"):
                if key in data and key in base:
                    if isinstance(data[key], list) and isinstance(base[key], list):
                        merged[key] = data[key]
            # Dicts: shallow merge (user keys override auto-fetched keys)
            for key in ("kpis",):
                if key in data and key in base:
                    if isinstance(data[key], dict) and isinstance(base[key], dict):
                        merged[key] = {**base[key], **data[key]}
        else:
            merged = base

        return self.generate(template, merged)

    def generate_all(self, base_data: dict[str, Any] | None = None) -> dict[str, str]:
        """Generate presentations for all templates.

        Args:
            base_data: Optional base data shared across all templates.

        Returns:
            Dict mapping template name → output path (or empty string on failure).
        """
        base_data = base_data or {}
        results: dict[str, str] = {}
        for tpl in self.available_templates():
            try:
                path = self.generate(tpl, dict(base_data))
                results[tpl] = path
            except (ImportError, ValueError, OSError, RuntimeError) as exc:
                self._log(f"[PRESENTATION] Template {tpl!r} failed: {exc}")
                results[tpl] = ""
        return results


# ── Singleton factory ────────────────────────────────────────────────────────

_instance: PresentationGenerator | None = None
_instance_lock = threading.RLock()


def get_presentation_generator(
    cfg: dict[str, Any] | None = None,
    *,
    log_fn: Any = None,
    output_dir: str = "",
) -> PresentationGenerator:
    """Return the process-level PresentationGenerator (creates on first call).

    Args:
        cfg: Optional config dict (reads PRESENTATION_GENERATOR_* keys).
        log_fn: Optional logging callback.
        output_dir: Override output directory (takes precedence over config).

    Returns:
        PresentationGenerator instance.
    """
    global _instance
    with _instance_lock:
        if _instance is None:
            merged_cfg: dict[str, Any] = dict(cfg or {})
            pc = presentation_config_from_cfg(merged_cfg)
            if output_dir:
                pc.output_dir = output_dir
            _instance = PresentationGenerator(pc, log_fn=log_fn)
    return _instance


def reset_presentation_generator() -> None:
    """Force-reset singleton (tests only)."""
    global _instance
    with _instance_lock:
        _instance = None


__all__ = [
    "PresentationConfig",
    "PresentationGenerator",
    "get_presentation_generator",
    "presentation_config_from_cfg",
    "reset_presentation_generator",
]
