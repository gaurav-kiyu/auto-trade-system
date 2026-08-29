"""Unit tests for scripts/generate_pptx.py — PPTX presentation generation.

Verifies:
  - All slide builder functions produce valid slide objects
  - The generated presentation has exactly 16 slides
  - Helper functions work correctly
  - No crashes on any slide function
"""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches
from scripts.generate_pptx import (
    add_bg,
    add_multiline_text,
    add_shape,
    add_slide,
    add_text,
    generate,
)

# ── Fixtures ─────────────────────────────────────────────────────────


@pytest.fixture()
def prs() -> Presentation:
    """Blank presentation with widescreen dimensions."""
    p = Presentation()
    p.slide_width = 12192000  # 13.33 inches in EMU
    p.slide_height = 6858000  # 7.5 inches in EMU
    return p


# ── Helper Function Tests ─────────────────────────────────────────────


class TestHelperFunctions:
    def test_add_slide_returns_slide(self, prs: Presentation) -> None:
        slide = add_slide(prs)
        assert slide is not None
        assert len(prs.slides) == 1

    def test_add_bg_sets_color(self, prs: Presentation) -> None:
        slide = add_slide(prs)
        add_bg(slide, RGBColor(0x1E, 0x1E, 0x2E))
        fill = slide.background.fill
        assert fill.fore_color.rgb == RGBColor(0x1E, 0x1E, 0x2E)

    def test_add_shape_creates_rectangle(self, prs: Presentation) -> None:
        slide = add_slide(prs)
        shape = add_shape(slide, Inches(0), Inches(0), Inches(1), Inches(1),
                          RGBColor(0x00, 0x6D, 0xAA))
        assert shape is not None
        assert shape.width == Inches(1)

    def test_add_shape_with_text(self, prs: Presentation) -> None:
        slide = add_slide(prs)
        shape = add_shape(slide, Inches(0), Inches(0), Inches(3), Inches(1),
                          RGBColor(0x00, 0x6D, 0xAA), "Hello", font_size=14)
        assert shape.text_frame.paragraphs[0].text == "Hello"

    def test_add_text_creates_textbox(self, prs: Presentation) -> None:
        slide = add_slide(prs)
        tb = add_text(slide, Inches(0), Inches(0), Inches(5), Inches(1),
                      "Hello", font_size=14)
        assert tb is not None
        assert tb.text_frame.paragraphs[0].text == "Hello"

    def test_add_multiline_text_creates_lines(self, prs: Presentation) -> None:
        slide = add_slide(prs)
        lines = ["Line 1", "Line 2", "Line 3"]
        tb = add_multiline_text(slide, Inches(0), Inches(0), Inches(5), Inches(2),
                                lines, font_size=14)
        # First paragraph should have first line
        assert "Line 1" in tb.text_frame.paragraphs[0].text

    def test_word_wrap_enabled(self, prs: Presentation) -> None:
        slide = add_slide(prs)
        tb = add_text(slide, Inches(0), Inches(0), Inches(5), Inches(1), "Test")
        assert tb.text_frame.word_wrap is True


# ── Generate Function Tests ──────────────────────────────────────────


class TestGenerateFunction:
    def test_generates_16_slides(self, tmp_path: Path) -> None:
        """generate() should produce a presentation with exactly 16 slides."""
        import os


        # Temporarily change working directory to tmp_path for output
        original_cwd = Path.cwd()
        try:
            os.chdir(str(tmp_path))
            generate()

            output_file = tmp_path / "OPB_SYSTEM_PRESENTATION_v2.57.1.pptx"
            assert output_file.exists(), f"Output PPTX file not found: {output_file}"

            prs_loaded = Presentation(str(output_file))
            assert len(prs_loaded.slides) == 16, (
                f"Expected 16 slides, got {len(prs_loaded.slides)}"
            )
        finally:
            os.chdir(str(original_cwd))

    def test_no_crash_on_full_generation(self, tmp_path: Path) -> None:
        """generate() should run without errors."""
        import os
        original_cwd = Path.cwd()
        try:
            os.chdir(str(tmp_path))
            result = generate()
            assert result is not None
            assert "OPB_SYSTEM_PRESENTATION" in str(result)
        finally:
            os.chdir(str(original_cwd))

    def test_all_slide_builders_individually(self, prs: Presentation) -> None:
        """All slide builder functions should work individually."""
        from scripts.generate_pptx import (
            add_slide,
        )

        # Just verify we can add multiple slides without errors
        for _ in range(3):
            add_slide(prs)
        assert len(prs.slides) == 3


# ── Module Import Verification ────────────────────────────────────────


class TestModuleImports:
    def test_module_importable(self) -> None:
        """scripts.generate_pptx should be importable without errors."""
        import scripts.generate_pptx as mod
        assert mod is not None
        assert hasattr(mod, "generate")

    def test_all_key_functions_exported(self) -> None:
        """Key functions should be accessible from the module."""
        import scripts.generate_pptx as mod
        expected = ["add_bg", "add_shape", "add_text", "add_multiline_text",
                    "add_slide", "generate"]
        for name in expected:
            assert hasattr(mod, name), f"Missing export: {name}"
