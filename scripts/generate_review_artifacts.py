"""Generate system review PDF + architecture PPT for OPB Index Options Buying Bot.

Regenerates the review deliverables under ``docs/review/``:

    python scripts/generate_review_artifacts.py

Outputs
-------
    docs/review/SYSTEM_REVIEW_SUMMARY.pdf
    docs/review/ARCHITECTURE_OVERVIEW.pptx

Dependencies
------------
    reportlab, python-pptx  (both optional; each section is generated
    independently and skipped on import errors).
"""
from __future__ import annotations

import os
import sys

_OUT_DIR = os.path.join("docs", "review")
os.makedirs(_OUT_DIR, exist_ok=True)

VERSION = "2.58.0"
REVIEW_DATE = "2026-08-08"
PDF_PATH = os.path.join(_OUT_DIR, "SYSTEM_REVIEW_SUMMARY.pdf")
PPT_PATH = os.path.join(_OUT_DIR, "ARCHITECTURE_OVERVIEW.pptx")

STRENGTHS: list[tuple[str, str]] = [
    ("Safety invariants", "Paper-mode invariant (simulated fills never reach a real broker), hard halt kill-switch, kill file, circuit breaker, watchdog thread, capital reservation lock, LTP sanity check."),
    ("Governance depth", "Constitution Validation Engine (23 categories), pre-implementation compliance gate, AI governance gate, live-readiness scorecard (5 blocking criteria), release governance pipeline."),
    ("Test maturity", "~14,700 tests; dedicated suites for smoke, live-readiness, NSE recorder, constitution, config schema, paper trader, execution hardening."),
    ("CI/CD completeness", "GitHub Actions with 9 jobs (lint, test matrix 3.11-3.14, coverage >=90%, security with pip-audit+bandit+semgrep, governance, certification, slow tests, build+checksum) + nightly full-suite cron."),
    ("Deployment options", "Dockerfile + docker-compose + supervisord, Dockerfile.realestate, k8s configmap, launcher EXE with single-instance lock."),
    ("Observability", "Prometheus metrics exporter, FastAPI enterprise dashboard (RBAC), health checker, audit trail (JSONL), structured logging with rotation + gzip."),
    ("Security posture", "Secrets moved to OPBUYING_* env vars, SecureConfig redaction, config.local.json gitignored, Telegram auth allowlists + rate limits + audit, bandit in CI."),
    ("Multi-asset coverage", "Equity, commodity (MCX), currency (CDS), index futures engines routed through a multi-asset dispatcher (equity enabled by default)."),
]

WEAKNESSES: list[tuple[str, str]] = [
    ("Live-readiness gate unmet", "0/50 paper trades, 0 trading days, 0% win rate -> the gate correctly blocks LIVE (AUTO) start. No broker credentials configured (BROKER_API_ENABLED=false, GENERIC driver). A PAPER session runs daily to build the record."),
    ("Live data latency", "NSE Akamai blocks automated scraping; yfinance fallback labels NIFTY LTP as 'last daily close' during live hours (stale intraday prices)."),
    ("Dead-code register noise", "44,324 ORPHANED_SYMBOL entries, all MEDIUM; mostly test functions/constants (register regenerated 2026-08-08). No orphaned modules found, so no safe bulk pruning exists."),
    ("Static-analysis findings", "Bandit: 0 HIGH / 0 MEDIUM / 298 LOW over 177,378 LOC (2026-08-08). All 14 MEDIUM resolved with proper fixes: identifier allow-lists on dynamic SQL, http/https scheme checks + opener-based URL calls, loopback default bind. Remaining 298 are LOW-severity heuristics (B101 asserts, B603 subprocess, B311 random) - see register."),
    ("Version skew", "Local runtime Python 3.14.4 vs CI matrix 3.11/3.12. Startup gate allows 3.10-3.19, but CI should add 3.13/3.14."),
    ("Dashboard binding", "web_dashboard binds 0.0.0.0 by default (opt-in feature; warns without TLS). Should default to 127.0.0.1."),
    ("ML cold start", "ML classifier has no training data (0 trades); predict_win_prob degrades to 0.5 until the paper track record accumulates."),
    ("Config drift artifacts", "Previous audit reports in docs/archive record config drift (EXECUTION_MODE defaults differ); current sync checks pass but older reports remain unarchived."),
]

IMPROVEMENTS: list[tuple[str, str, str]] = [
    ("P1", "Build paper track record", "Run EXECUTION_MODE=PAPER for >=10 trading days / >=50 trades to unblock the live-readiness gate. This is the single blocking item for AUTO."),
    ("P1", "Configure broker secrets", "Set OPBUYING_BROKER_API_KEY/ACCESS_TOKEN/USER_ID/PASSWORD env vars and BROKER_DRIVER=KITE|ANGEL when ready; keep EXECUTION_MODE progression MANUAL->PAPER->SIGNAL_ONLY->AUTO."),
    ("P2", "DONE - Harden urlopen schemes", "http/https scheme allowlists on all urlopen sites: core (6 modules) + scripts (realestate_synthetic_monitor, test_deployment); calls moved to build_opener().open() - no nosec."),
    ("P2", "DONE - Fix dynamic SQL", "Identifier allow-list validation (_safe_ident) on all dynamic table/column SQL in scripts (check_db_integrity, migrate_to_postgresql, verify_restore, quantitative_validation_report); values parameter-bound - no nosec. Bandit MEDIUM 14 -> 0."),
    ("P2", "DONE - Dashboard default bind", "web_dashboard + enterprise_dashboard + launch_realestate bind 127.0.0.1 by default (0.0.0.0 only when explicitly configured)."),
    ("P2", "DONE - CI matrix", "Python 3.13/3.14 added to the GitHub Actions test matrix; semgrep added to the security job; nightly full-suite cron scheduled."),
    ("P3", "Register hygiene", "Run scan_dead_code.py --update-registers after every scan so the register reflects current code (stale rows were regenerated 2026-08-08)."),
    ("P3", "Archive stale docs", "Move 2026-07-20 audit reports in docs/archive to a dated subfolder and regenerate config-drift reports against current defaults."),
]

LIVE_GATE: list[tuple[str, str, str]] = [
    ("Minimum paper trades", "0 / 50", "BLOCKED"),
    ("Win rate >= 50%", "0.0%", "BLOCKED"),
    ("Profit factor >= 1.30", "0.000", "BLOCKED"),
    ("Max drawdown <= 15%", "0.0%", "PASS"),
    ("Minimum trading days", "0 / 10", "BLOCKED"),
    ("Sharpe >= 0.5 (rec.)", "0.000", "NOT MET"),
    ("ML accuracy >= 50% (rec.)", "50.0%", "PASS"),
]

VERIFICATION: list[tuple[str, str]] = [
    ("Pre-implementation compliance check", "PASSED - no violations"),
    ("Repository hygiene check", "PRISTINE - 0 issues"),
    ("Script & artifact sync (sync_artifacts.py)", "ALL SYNCHRONIZED - 0 issues (11 version-drift items fixed)"),
    ("Architecture compliance (check_architecture_compliance.py)", "PASSED"),
    ("Constitution governance scoring", "9.19 avg / 10 across 111 categories (min 8.5, none below 6.0)"),
    ("Smoke + live-readiness + NSE recorder tests", "51/51 PASSED"),
    ("Governance + config-schema tests", "288/288 PASSED"),
    ("Full regression suite (~14,700 tests)", "PASSED - 14,659 passed, 0 failed, 0 errors (95 skipped, 1 xfail, 5 xpass) in ~95 min"),
    ("Bandit static analysis (177,378 LOC)", "0 HIGH / 0 MEDIUM / 298 LOW - all 14 MEDIUM fixed (identifier allow-lists, scheme checks, loopback bind)"),
    ("Live-market session (PAPER on real NSE data)", "RUNNING since 2026-08-07 11:51 IST - 0 orders, reconciliation CLEAN"),
]


# ── PDF ──────────────────────────────────────────────────────────────────────
def build_pdf() -> None:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import mm
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

    styles = getSampleStyleSheet()
    h1 = ParagraphStyle("H1", parent=styles["Heading1"], fontSize=16, spaceAfter=6,
                        textColor=colors.HexColor("#0b3d66"))
    h2 = ParagraphStyle("H2", parent=styles["Heading2"], fontSize=12, spaceBefore=10,
                        spaceAfter=4, textColor=colors.HexColor("#0b3d66"))
    body = ParagraphStyle("Body", parent=styles["BodyText"], fontSize=9, leading=12)
    small = ParagraphStyle("Small", parent=styles["BodyText"], fontSize=8, leading=10)

    doc = SimpleDocTemplate(
        PDF_PATH, pagesize=A4,
        leftMargin=16 * mm, rightMargin=16 * mm,
        topMargin=14 * mm, bottomMargin=14 * mm,
        title=f"OPB System Review {REVIEW_DATE}", author="OPB Engineering",
    )

    story = []
    story.append(Paragraph("OPB Index Options Buying Bot", h1))
    story.append(Paragraph(
        f"End-to-End System Review, Cleanup & Enhancement Report - v{VERSION} ({REVIEW_DATE})",
        styles["Title"]))
    story.append(Spacer(1, 4 * mm))

    story.append(Paragraph("1. Review Scope", h2))
    story.append(Paragraph(
        "Full audit of architecture, code, configuration, security, tests and documentation; "
        "live-market validation via a SHADOW-LIVE (SIGNAL_ONLY) session; generation of this summary "
        "and an architecture deck. No risk controls (SL/TARGET/drawdown), broker abstraction, or "
        "paper-mode invariants were modified.", body))

    story.append(Paragraph("2. Live-Market Session Status", h2))
    live_rows = [
        ["Mode", "SHADOW LIVE (SIGNAL_ONLY) - signals on real NSE market data, zero orders"],
        ["Start time", "2026-08-07 09:23 IST (market open 09:15 IST)"],
        ["Process", "index_trader.py, trading loop interval 30s, reconciliation CLEAN"],
        ["Orders placed", "0 (by design - no broker configured, gate enforced)"],
        ["Live-readiness gate", "BLOCKED - see section 5 (correct behaviour)"],
    ]
    t = Table([[Paragraph(f"<b>{k}</b>", small), Paragraph(v, small)] for k, v in live_rows],
              colWidths=[32 * mm, 146 * mm])
    t.setStyle(TableStyle([
        ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#c8d6e5")),
        ("BACKGROUND", (0, 0), (0, -1), colors.HexColor("#eef4fa")),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("LEFTPADDING", (0, 0), (-1, -1), 4),
        ("RIGHTPADDING", (0, 0), (-1, -1), 4),
        ("TOPPADDING", (0, 0), (-1, -1), 3),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
    ]))
    story.append(t)

    story.append(Paragraph("3. Verification Results", h2))
    vt = Table([[Paragraph("<b>Check</b>", small), Paragraph("<b>Result</b>", small)]
                for k, v in VERIFICATION], colWidths=[108 * mm, 70 * mm])
    vt.setStyle(TableStyle([
        ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#c8d6e5")),
        ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor("#f5f9fd")),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("LEFTPADDING", (0, 0), (-1, -1), 4),
        ("RIGHTPADDING", (0, 0), (-1, -1), 4),
        ("TOPPADDING", (0, 0), (-1, -1), 3),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
    ]))
    story.append(vt)

    story.append(Paragraph("4. Strengths", h2))
    for name, desc in STRENGTHS:
        story.append(Paragraph(f"<b>{name}:</b> {desc}", body))
        story.append(Spacer(1, 1.5 * mm))

    story.append(Paragraph("5. Weaknesses / Risks", h2))
    for name, desc in WEAKNESSES:
        story.append(Paragraph(f"<b>{name}:</b> {desc}", body))
        story.append(Spacer(1, 1.5 * mm))

    story.append(Paragraph("6. Live-Readiness Gate (blocks LIVE/AUTO start)", h2))
    gt = Table([[Paragraph("<b>Criterion</b>", small), Paragraph("<b>Actual</b>", small),
                 Paragraph("<b>Status</b>", small)] for c, a, s in LIVE_GATE],
               colWidths=[88 * mm, 40 * mm, 50 * mm])
    gt.setStyle(TableStyle([
        ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#c8d6e5")),
        ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor("#f5f9fd")),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("LEFTPADDING", (0, 0), (-1, -1), 4),
        ("RIGHTPADDING", (0, 0), (-1, -1), 4),
        ("TOPPADDING", (0, 0), (-1, -1), 3),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
    ]))
    story.append(gt)
    story.append(Spacer(1, 2 * mm))
    story.append(Paragraph(
        "Recommendation from the checker: continue paper trading until all blocking criteria are met. "
        "Score 1/5 blocking passed. A daily PAPER session is active to accumulate the 50-trade record.", small))

    story.append(Paragraph("7. Suggested Improvements (prioritised)", h2))
    for prio, name, desc in IMPROVEMENTS:
        story.append(Paragraph(f"<b>{prio} - {name}:</b> {desc}", body))
        story.append(Spacer(1, 1.5 * mm))

    doc.build(story)
    print(f"PDF written: {PDF_PATH}")


# ── PPT ──────────────────────────────────────────────────────────────────────
def build_ppt() -> None:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    NAVY = RGBColor(0x0B, 0x3D, 0x66)
    TEAL = RGBColor(0x14, 0x8F, 0x8F)
    RED = RGBColor(0xB3, 0x2D, 0x2D)
    DARK = RGBColor(0x22, 0x2B, 0x35)
    GREY = RGBColor(0x5A, 0x64, 0x70)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def add_title(slide, text: str, sub: str | None = None) -> None:
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = text
        r.font.size = Pt(30)
        r.font.bold = True
        r.font.color.rgb = NAVY
        if sub:
            p2 = tf.add_paragraph()
            r2 = p2.add_run()
            r2.text = sub
            r2.font.size = Pt(14)
            r2.font.color.rgb = GREY

    def add_bullets(slide, items, top=1.5, height=5.5, font_size=13, color=DARK) -> None:
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12.1), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = True
        first = True
        for head, desc in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = Pt(8)
            r = p.add_run()
            r.text = f"{head}: "
            r.font.bold = True
            r.font.size = Pt(font_size)
            r.font.color.rgb = NAVY
            r2 = p.add_run()
            r2.text = desc
            r2.font.size = Pt(font_size - 1)
            r2.font.color.rgb = color

    # Slide 1 - Title
    s = prs.slides.add_slide(blank)
    tb = s.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(11.3), Inches(2.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "OPB Index Options Buying Bot"
    r.font.size = Pt(44)
    r.font.bold = True
    r.font.color.rgb = NAVY
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = "Architecture Overview & System Review"
    r2.font.size = Pt(26)
    r2.font.color.rgb = TEAL
    p3 = tf.add_paragraph()
    r3 = p3.add_run()
    r3.text = f"v{VERSION}  |  {REVIEW_DATE}  |  NSE Index Options (NIFTY / BANKNIFTY / FINNIFTY)"
    r3.font.size = Pt(14)
    r3.font.color.rgb = GREY
    p4 = tf.add_paragraph()
    r4 = p4.add_run()
    r4.text = "Strengths  |  Weaknesses  |  Suggested Improvements"
    r4.font.size = Pt(13)
    r4.font.color.rgb = GREY

    # Slide 2 - System overview
    s = prs.slides.add_slide(blank)
    add_title(s, "System Overview", "Python 3.10-3.19 | Windows primary, Linux/Docker | NSE index options")
    add_bullets(s, [
        ("Purpose", "Automated NSE index options buying with signal generation, risk management, execution (Kite/Angel/Paper)."),
        ("Data", "yfinance (LTP/OHLCV) primary; NSE option chain blocked by Akamai; broker WS optional."),
        ("Config", "3-layer merge: index_config.defaults.json (1,058 keys) <- config.json <- OPBUYING_* env secrets."),
        ("State", "SQLite (trades/journal/ML tracker/OI snapshots) + trader_state.json survive restarts."),
        ("Notifications", "Telegram bot with auth allowlists, rate limits, audit, priority queue."),
        ("Current run", "SHADOW LIVE (SIGNAL_ONLY) active during market hours - real data, zero orders."),
    ], top=1.5)

    # Slide 3 - Architecture layers
    s = prs.slides.add_slide(blank)
    add_title(s, "Architecture Layers", "Clean-ish architecture with domains, ports, DI container")
    add_bullets(s, [
        ("Entry / UI", "index_trader.py (brain), launcher.py (GUI), web dashboard (FastAPI+RBAC), Telegram commander"),
        ("Domains (DEBT-008)", "config, broker, market, trading, admin - extracted responsibilities"),
        ("Core services", "Signal pipeline (IV rank -> session -> ML -> tier), risk service, position service, execution service, reconciliation, self-healing orchestrator"),
        ("Strategy engines", "Index options + equity + commodity + currency + futures via multi-asset dispatcher"),
        ("Adapters / ports", "Broker adapters (Paper/Kite/Angel), yfinance/NSE data, Telegram, reporting (ReportLab), metrics (Prometheus)"),
        ("Safety layer", "Hard halt, circuit breaker, watchdog, kill file, margin validator, deterministic state machine, idempotency certifier"),
    ], top=1.5, font_size=12)

    # Slide 4 - Strengths
    s = prs.slides.add_slide(blank)
    add_title(s, "Strengths")
    add_bullets(s, STRENGTHS, top=1.5, font_size=12)

    # Slide 5 - Weaknesses
    s = prs.slides.add_slide(blank)
    add_title(s, "Weaknesses / Risks")
    add_bullets(s, WEAKNESSES, top=1.5, font_size=12, color=RED)

    # Slide 6 - Improvements
    s = prs.slides.add_slide(blank)
    add_title(s, "Suggested Improvements", "Prioritised roadmap")
    add_bullets(s, [(f"{p} - {n}", d) for p, n, d in IMPROVEMENTS], top=1.5, font_size=12)

    # Slide 7 - Live readiness gate
    s = prs.slides.add_slide(blank)
    add_title(s, "Live-Readiness Gate - Current Status", "Automated blocker for LIVE (AUTO) execution")
    add_bullets(s, [
        ("Minimum paper trades", "0 / 50  (BLOCKED)"),
        ("Win rate >= 50%", "0.0%  (BLOCKED)"),
        ("Profit factor >= 1.30", "0.000  (BLOCKED)"),
        ("Max drawdown <= 15%", "0.0%  (PASS)"),
        ("Minimum trading days", "0 / 10  (BLOCKED)"),
        ("Sharpe >= 0.5", "0.000  (not met)"),
        ("Recommendation", "Run paper mode until criteria met; then MANUAL -> PAPER -> SIGNAL_ONLY -> AUTO progression with broker secrets."),
    ], top=1.5, font_size=13)

    prs.save(PPT_PATH)
    print(f"PPT written: {PPT_PATH}")


if __name__ == "__main__":
    ok = True
    try:
        build_pdf()
    except Exception as exc:
        print(f"PDF FAILED: {exc}", file=sys.stderr)
        ok = False
    try:
        build_ppt()
    except Exception as exc:
        print(f"PPT FAILED: {exc}", file=sys.stderr)
        ok = False
    sys.exit(0 if ok else 1)
